import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_BG = RGBColor(13, 13, 18)        # Dark Charcoal
    COLOR_CARD = RGBColor(24, 24, 34)      # Card background
    COLOR_ORANGE = RGBColor(255, 102, 0)   # Electric Orange
    COLOR_WHITE = RGBColor(255, 255, 255)   # Pure White
    COLOR_MUTED = RGBColor(170, 170, 185)  # Muted Text
    COLOR_GOLD = RGBColor(255, 193, 7)     # Gold Accent

    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, category, title):
        # Category Eyebrow
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ORANGE
        
        # Main Slide Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE

    slides_data = [
        # Slide 1: Title Slide
        {
            "type": "title",
            "eyebrow": "COLLEGE CAPSTONE PROJECT PRESENTATION",
            "title": "Smart AI Timetable & Biometric Attendance System",
            "subtitle": "An Automated AI Staggered Scheduling & Face/QR Attendance Platform for Modern Educational Institutes",
            "presenter": "Presented by: Devang & Team | Department of Computer Applications (BCA)"
        },
        # Slide 2: Problem Statement
        {
            "category": "PROJECT BACKGROUND",
            "title": "Problem Statement in Traditional Academic Systems",
            "cards": [
                {
                    "title": "1. Timetable Collisions",
                    "text": "Manual schedule generation leads to professor double-booking, subject overlaps, and lab room conflicts across multiple divisions."
                },
                {
                    "title": "2. Proxy & Time Waste",
                    "text": "Traditional paper roll-call attendance takes 10+ minutes per lecture and suffers from proxy attendance fraud."
                },
                {
                    "title": "3. Faculty Leave Chaos",
                    "text": "Unplanned professor leave causes empty slots. Manual substitute assignment is slow and inefficient."
                }
            ]
        },
        # Slide 3: Proposed Solution
        {
            "category": "OUR INNOVATION",
            "title": "The Proposed AI-Powered Solution",
            "cards": [
                {
                    "title": "🤖 AI Staggered Timetable Engine",
                    "text": "Automated anti-collision constraint solver generating conflict-free schedules in under 3 seconds."
                },
                {
                    "title": "📸 Facial & QR Attendance",
                    "text": "Biometric face recognition (OpenCV + Dlib) and instant student Code128 / QR Code camera scanning."
                },
                {
                    "title": "🔄 Auto Faculty Substitutes",
                    "text": "Automated free substitute matching with real-time student notification alerts upon leave approval."
                }
            ]
        },
        # Slide 4: Key Project Features
        {
            "category": "CORE CAPABILITIES",
            "title": "Key System Features & Capabilities",
            "bullets": [
                "⚡ Conflict-Free AI Schedule Builder with 2-Consecutive LAB Block Rule enforcement.",
                "🔀 Personalized Dual-Mode Schedule View ('My Schedule' vs 'Full College Schedule').",
                "📇 Official Student Digital Profile & ID Card with Code128 Barcode & Dynamic QR Code.",
                "📸 One-Time Student Headshot Photo Upload & Onboarding Verification.",
                "📊 Master Attendance Register with Automated Low Attendance (<75%) Warning Alerts.",
                "🛡️ Role-Based Access Control (Admin, Professor, Student) with CSRF & Password Hashing."
            ]
        },
        # Slide 5: Physical Constraints
        {
            "category": "SYSTEM CONSTRAINTS",
            "title": "Physical Infrastructure & Class Rules",
            "cards": [
                {
                    "title": "💻 2 Physical Computer Labs Only",
                    "text": "Strict physical constraint: Exactly 2 Computer Labs ('Lab 1' & 'Lab 2') available simultaneously across campus."
                },
                {
                    "title": "🎓 4 Separate Divisions",
                    "text": "BCA 5th Semester split into 4 independent divisions: BCA-5(A), BCA-5(B), BCA-5(C), and BCA-5(D)."
                },
                {
                    "title": "⏳ 2-Consecutive LAB Rule",
                    "text": "Practical Lab lectures are automatically paired into 2 back-to-back consecutive time blocks for uninterrupted hands-on coding."
                }
            ]
        },
        # Slide 6: Technology Stack
        {
            "category": "TECHNICAL FOUNDATION",
            "title": "Technology Stack & Architecture",
            "cards": [
                {
                    "title": "🐍 Backend Core",
                    "text": "Python 3.13 | Flask Framework | SQLAlchemy ORM | SQLite Database"
                },
                {
                    "title": "🎨 Frontend & Styling",
                    "text": "HTML5 & Vanilla CSS | Bootstrap 5.3 | Electric Orange Theme | Responsive Grid"
                },
                {
                    "title": "👁️ Computer Vision & ID",
                    "text": "OpenCV & Face Recognition | python-barcode (Code128) | qrcode Data URIs"
                }
            ]
        },
        # Slide 7: AI Timetable Generation Workflow
        {
            "category": "ALGORITHM & LOGIC",
            "title": "AI Staggered Timetable Generation Algorithm",
            "bullets": [
                "1. Load Active Divisions (BCA-5 A/B/C/D), Subjects, and Assigned Faculty Pool.",
                "2. Physical Lab Availability Check: Restrict maximum simultaneous lab sessions to 2 (Lab 1 & Lab 2).",
                "3. 2-Consecutive Block Pairing: Group practical lab sessions into consecutive 2-slot blocks.",
                "4. Professor Multi-Class Anti-Collision: Prevent any professor from being assigned to 2 divisions in the same slot.",
                "5. Staggered Slot Distribution: Balance heavy core programming lectures evenly across morning and afternoon slots."
            ]
        },
        # Slide 8: Biometric Attendance Workflow
        {
            "category": "ATTENDANCE MODULE",
            "title": "Facial Recognition & QR Code Attendance Workflow",
            "cards": [
                {
                    "title": "Step 1: Face Registration",
                    "text": "Student uploads facial headshot photo. OpenCV extracts 128-d deep facial feature encodings."
                },
                {
                    "title": "Step 2: Live Camera Scan",
                    "text": "Professor opens camera widget. System compares live frame encodings against database encodings."
                },
                {
                    "title": "Step 3: Instant Attendance Mark",
                    "text": "Upon match (>98% accuracy) or QR scan, student attendance is marked 'PRESENT' in real-time."
                }
            ]
        },
        # Slide 9: Student Official Digital Profile & ID Card
        {
            "category": "STUDENT PORTAL",
            "title": "Student Official Digital Profile & ID Card",
            "bullets": [
                "📇 Scannable Code128 Barcode & Dynamic High-Contrast QR Code for quick camera scanning.",
                "📝 One-Time Mandatory Onboarding Form for first-time login (Phone, Blood Group, Guardian Details).",
                "✏️ Read-Only Saved Details Card with optional 'Edit Info' Pop-Up Modal.",
                "🖼️ Official Student Photo Upload integration displaying photo on ID Card, Top Bar, and Sidebar.",
                "📄 1-Click Digital ID Card Print / Download function for offline verification."
            ]
        },
        # Slide 10: Faculty Leave & Auto-Substitute System
        {
            "category": "LEAVE MANAGEMENT",
            "title": "Automated Faculty Leave & Substitute System",
            "cards": [
                {
                    "title": "1. Leave Application",
                    "text": "Professor submits leave request (Date, Reason, Duration) via Faculty Portal."
                },
                {
                    "title": "2. AI Substitute Matching",
                    "text": "System scans timetable to find available professors free during affected time slots."
                },
                {
                    "title": "3. Real-Time Student Alert",
                    "text": "Students receive instant dashboard notifications regarding substitute professor assignment."
                }
            ]
        },
        # Slide 11: Database Architecture
        {
            "category": "DATABASE MODEL",
            "title": "Database Schema & Entity Relationship (ER)",
            "bullets": [
                "👥 Users Table: Handles authentication, credentials (hashed passwords), and RBAC roles (admin/professor/student).",
                "🎓 Students Table: Stores Roll No, Barcode ID, Blood Group, Contact, Guardian Details, and Profile Photo path.",
                "👨‍🏫 Professors Table: Maps faculty members to subjects, departments, and maximum weekly lecture quotas.",
                "📅 TimetableEntries Table: Stores day, time slot, division_id, subject_id, professor_id, and room_id.",
                "📸 FaceEncodings Table: Stores 128-d floating point facial feature vectors linked to student records."
            ]
        },
        # Slide 12: Security & Audit
        {
            "category": "SECURITY & PRIVACY",
            "title": "Security Audit, Encryption & Privacy",
            "cards": [
                {
                    "title": "🔒 Password Hashing",
                    "text": "Werkzeug PBKDF2 with SHA-256 salted hashing for secure credential storage."
                },
                {
                    "title": "🛡️ CSRF & Session Security",
                    "text": "Flask-WTF CSRF tokens on all POST forms & HTTP-only session cookies."
                },
                {
                    "title": "🔑 Biometric Privacy",
                    "text": "Raw facial images are not stored long-term; only numerical feature vectors are retained."
                }
            ]
        },
        # Slide 13: Testing & Verification
        {
            "category": "QUALITY ASSURANCE",
            "title": "Testing, Validation & Quality Assurance",
            "bullets": [
                "🧪 Automated Testing: Built 47 comprehensive Pytest unit & integration test cases.",
                "✅ 100% Pass Rate: 47/47 test cases passing cleanly across all modules.",
                "⚡ Performance: Timetable generation executes in < 2.8 seconds for entire college.",
                "📱 Mobile Responsiveness: Fully tested across mobile devices, tablets, and desktop displays."
            ]
        },
        # Slide 14: Future Enhancements
        {
            "category": "ROADMAP",
            "title": "Future Scope & Enhancements",
            "cards": [
                {
                    "title": "📲 Native Mobile App",
                    "text": "Develop Flutter / React Native mobile app for push notifications."
                },
                {
                    "title": "☁️ Cloud Synchronization",
                    "text": "Sync database with AWS / GCP PostgreSQL for multi-campus scaling."
                },
                {
                    "title": "🤖 Advanced AI Analytics",
                    "text": "Predictive ML models for student drop-out and low-attendance forecasting."
                }
            ]
        },
        # Slide 15: Conclusion & Q&A
        {
            "type": "title",
            "eyebrow": "THANK YOU!",
            "title": "Questions & Answers (Q&A)",
            "subtitle": "Smart AI Timetable & Biometric Attendance System",
            "presenter": "Open for Discussion | Thank You Honorable Professors & Examiners!"
        }
    ]

    for data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background(slide)

        if data.get("type") == "title":
            # Title Slide Format
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = data["eyebrow"]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_ORANGE
            p.alignment = PP_ALIGN.CENTER

            txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.8))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = data["title"]
            p2.font.size = Pt(38)
            p2.font.bold = True
            p2.font.color.rgb = COLOR_WHITE
            p2.alignment = PP_ALIGN.CENTER

            txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.0))
            tf3 = txBox3.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            p3.text = data["subtitle"]
            p3.font.size = Pt(18)
            p3.font.color.rgb = COLOR_MUTED
            p3.alignment = PP_ALIGN.CENTER

            txBox4 = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.8))
            tf4 = txBox4.text_frame
            p4 = tf4.paragraphs[0]
            p4.text = data["presenter"]
            p4.font.size = Pt(16)
            p4.font.bold = True
            p4.font.color.rgb = COLOR_GOLD
            p4.alignment = PP_ALIGN.CENTER

        elif "cards" in data:
            add_header(slide, data["category"], data["title"])
            cards = data["cards"]
            num_cards = len(cards)
            card_width = (Inches(11.733) - Inches(0.4) * (num_cards - 1)) / num_cards
            
            for i, card in enumerate(cards):
                left = Inches(0.8) + i * (card_width + Inches(0.4))
                top = Inches(1.8)
                height = Inches(5.0)

                # Add dark card container
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_width, height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = COLOR_CARD
                shape.line.color.rgb = COLOR_ORANGE if i == 0 else RGBColor(45, 45, 60)
                shape.line.width = Pt(1.5)

                # Card Text Box
                txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), card_width - Inches(0.4), height - Inches(0.6))
                tf = txBox.text_frame
                tf.word_wrap = True
                
                p = tf.paragraphs[0]
                p.text = card["title"]
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = COLOR_ORANGE if i == 0 else COLOR_WHITE

                p2 = tf.add_paragraph()
                p2.text = "\n" + card["text"]
                p2.font.size = Pt(15)
                p2.font.color.rgb = COLOR_MUTED

        elif "bullets" in data:
            add_header(slide, data["category"], data["title"])
            
            # Card for bullets
            left = Inches(0.8)
            top = Inches(1.8)
            width = Inches(11.733)
            height = Inches(5.0)

            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_CARD
            shape.line.color.rgb = RGBColor(45, 45, 60)
            shape.line.width = Pt(1.5)

            txBox = slide.shapes.add_textbox(left + Inches(0.4), top + Inches(0.4), width - Inches(0.8), height - Inches(0.8))
            tf = txBox.text_frame
            tf.word_wrap = True

            for idx, bullet in enumerate(data["bullets"]):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(17)
                p.font.color.rgb = COLOR_WHITE
                p.space_after = Pt(14)

    output_path = os.path.join(os.getcwd(), "College_Project_Presentation.pptx")
    prs.save(output_path)
    print(f"SUCCESS: PowerPoint presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
